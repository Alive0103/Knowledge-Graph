from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"D:\work\毕设\知识图谱\Knowledge-Graph")
OUT = ROOT / "王越洋-中期答辩PPT.pptx"

IMG_PREPROCESS = Path(r"C:\Users\alive\AppData\Roaming\Typora\typora-user-images\image-20260313131606367.png")
IMG_LINKING = Path(r"C:\Users\alive\AppData\Roaming\Typora\typora-user-images\image-20260313131738574.png")
IMG_ALIGN = Path(r"C:\Users\alive\AppData\Roaming\Typora\typora-user-images\image-20260313131752702.png")

TITLE = "融合内外部数据的知识百科构建与检索增强生成系统"
SUBTITLE = "本科毕业设计中期答辩"
AUTHOR = "王越洋"
STUDENT_ID = "22009200894"
MAJOR = "计算机科学与技术"

BG = RGBColor(246, 243, 238)
PRIMARY = RGBColor(31, 58, 95)
ACCENT = RGBColor(200, 106, 43)
GREEN = RGBColor(94, 139, 91)
LIGHT_BLUE = RGBColor(234, 241, 247)
LIGHT_ORANGE = RGBColor(248, 239, 231)
LIGHT_GREEN = RGBColor(238, 248, 238)
TEXT = RGBColor(38, 38, 38)
MUTED = RGBColor(103, 112, 124)
WHITE = RGBColor(255, 255, 255)


def set_background(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_top_bar(slide, title_tag="MIDTERM DEFENSE"):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.28))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.42), Inches(1.75), Inches(0.36))
    tag.fill.solid()
    tag.fill.fore_color.rgb = LIGHT_ORANGE
    tag.line.color.rgb = ACCENT
    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_tag
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_slide_title(slide, title, subtitle=None):
    add_top_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.48), Inches(0.95), Inches(9.8), Inches(0.65))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = PRIMARY
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(10), Inches(0.35))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = MUTED


def add_textbox(slide, left, top, width, height, text, font_size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return box


def add_bullets(slide, left, top, width, height, items, font_size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def add_panel(slide, left, top, width, height, fill=WHITE, line=RGBColor(220, 224, 230), radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_metric_card(slide, left, top, width, height, label, value, accent=PRIMARY):
    panel = add_panel(slide, left, top, width, height, fill=WHITE, line=RGBColor(224, 228, 234))
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.10))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    add_textbox(slide, left + Inches(0.18), top + Inches(0.18), width - Inches(0.3), Inches(0.28), label, font_size=11, color=MUTED, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.48), width - Inches(0.3), Inches(0.42), value, font_size=22, color=accent, bold=True)
    return panel


def fit_image(slide, image_path, left, top, width, height):
    image_path = Path(image_path)
    if not image_path.exists():
        return None
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    box_w = width / 914400
    box_h = height / 914400
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    if img_ratio > box_ratio:
        final_w = width
        final_h = int(width / img_ratio)
        final_left = left
        final_top = top + int((height - final_h) / 2)
    else:
        final_h = height
        final_w = int(height * img_ratio)
        final_top = top
        final_left = left + int((width - final_w) / 2)
    return slide.shapes.add_picture(str(image_path), final_left, final_top, width=final_w, height=final_h)


def add_flow_box(slide, left, top, width, height, title, body, fill, line_color, title_color=PRIMARY):
    add_panel(slide, left, top, width, height, fill=fill, line=line_color)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.14), width - Inches(0.32), Inches(0.32), title, font_size=14, color=title_color, bold=True)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.50), width - Inches(0.32), height - Inches(0.62), body, font_size=11.5, color=TEXT)


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.25)
    line.line.end_arrowhead = True
    return line


def add_chart_mrr(slide, left, top, width, height):
    chart_data = CategoryChartData()
    chart_data.categories = ["纯向量", "纯文本检索", "纯LLM", "向量+LLM重排", "动态混合"]
    chart_data.add_series("MRR", (0.0590, 0.6925, 0.7015, 0.0686, 0.6892))

    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 0.8
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.major_unit = 0.1
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "实体链接方案对比（全部数据，MRR）"
    chart.chart_title.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    series = chart.series[0]
    fill_colors = [RGBColor(196, 208, 224), PRIMARY, ACCENT, RGBColor(210, 214, 220), GREEN]
    for idx, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = fill_colors[idx]
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.font.size = Pt(10)


def add_small_table(slide, left, top, width, height, headers, rows, header_fill=PRIMARY):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    col_width = int(width / len(headers))
    for i in range(len(headers)):
        table.columns[i].width = col_width
    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Microsoft YaHei"
                r.font.bold = True
                r.font.size = Pt(11)
                r.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else RGBColor(248, 250, 252)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(11)
                    run.font.color.rgb = TEXT
                p.alignment = PP_ALIGN.CENTER
    return table_shape


def add_timeline(slide, left, top, width, bar_height):
    months = [
        ("3月下旬-4月中旬", "实体链接优化\n跨语言对齐补实验", ACCENT, 0.0, 0.46),
        ("4月中下旬", "Graph+RAG实验\n系统联调", PRIMARY, 0.49, 0.29),
        ("5月", "论文撰写修改\n图表整理与定稿", GREEN, 0.80, 0.20),
    ]
    axis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top + Inches(0.75), width, Pt(4))
    axis.fill.solid()
    axis.fill.fore_color.rgb = RGBColor(195, 201, 211)
    axis.line.fill.background()
    for idx, (label, task, color, start, ratio) in enumerate(months):
        seg_left = left + int(width * start)
        seg_w = int(width * ratio)
        seg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, seg_left, top + Inches(0.48), seg_w, bar_height)
        seg.fill.solid()
        seg.fill.fore_color.rgb = color
        seg.line.fill.background()
        add_textbox(slide, seg_left + Inches(0.08), top + Inches(0.14), seg_w - Inches(0.16), Inches(0.28), label, font_size=11, color=color, bold=True)
        body = slide.shapes.add_textbox(seg_left + Inches(0.14), top + Inches(0.72), seg_w - Inches(0.28), Inches(0.55))
        tf = body.text_frame
        tf.word_wrap = True
        tf.clear()
        lines = task.split("\n")
        for l_idx, line in enumerate(lines):
            p = tf.paragraphs[0] if l_idx == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(11.5)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(247, 242, 235))
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.35), 0, Inches(4.983), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PRIMARY
    panel.line.fill.background()
    accent_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.05), 0, Inches(0.18), Inches(7.5))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT
    accent_line.line.fill.background()

    add_textbox(slide, Inches(0.75), Inches(0.75), Inches(7.0), Inches(0.45), SUBTITLE, font_size=16, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.75), Inches(1.35), Inches(6.9), Inches(2.2), TITLE, font_size=28, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(0.78), Inches(3.95), Inches(6.5), Inches(0.5), "项目关键词：离线知识库 / 实体链接 / 跨语言实体对齐 / Graph+RAG", font_size=14, color=MUTED)

    info_panel = add_panel(slide, Inches(0.75), Inches(4.75), Inches(4.6), Inches(1.45), fill=WHITE, line=RGBColor(228, 226, 220))
    add_textbox(slide, Inches(0.98), Inches(5.02), Inches(4.2), Inches(0.3), f"答辩人：{AUTHOR}", font_size=16, color=TEXT, bold=True)
    add_textbox(slide, Inches(0.98), Inches(5.35), Inches(4.2), Inches(0.3), f"学号：{STUDENT_ID}", font_size=14, color=MUTED)
    add_textbox(slide, Inches(0.98), Inches(5.62), Inches(4.2), Inches(0.3), f"专业：{MAJOR}", font_size=14, color=MUTED)

    add_textbox(slide, Inches(8.68), Inches(1.15), Inches(3.8), Inches(0.6), "Knowledge\nGraph", font_size=24, color=WHITE, bold=True)
    add_textbox(slide, Inches(8.72), Inches(2.05), Inches(3.7), Inches(0.7), "融合内外部数据\n构建可检索、可对齐、可增强的知识底座", font_size=15, color=RGBColor(219, 228, 240))
    add_metric_card(slide, Inches(8.72), Inches(3.25), Inches(1.45), Inches(1.0), "实体规模", "20,000+", accent=ACCENT)
    add_metric_card(slide, Inches(10.27), Inches(3.25), Inches(1.45), Inches(1.0), "NER F1", "98.26%", accent=GREEN)
    add_metric_card(slide, Inches(11.82), Inches(3.25), Inches(1.15), Inches(1.0), "对齐MRR", "0.76", accent=RGBColor(123, 104, 171))

    # Slide 2: Background & objectives
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "研究背景与目标", "为什么做、要解决什么问题")
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.95), Inches(4.55), fill=WHITE)
    add_textbox(slide, Inches(0.8), Inches(2.25), Inches(5.4), Inches(0.4), "研究背景", font_size=18, color=PRIMARY, bold=True)
    add_bullets(
        slide, Inches(0.82), Inches(2.75), Inches(5.2), Inches(3.5),
        [
            "组织内部知识与外部开放知识分散、异构，缺少统一组织方式。",
            "军事领域实体命名复杂，简称、别名和型号表达多，直接检索容易误配。",
            "传统文本检索难以支撑可信、可溯源的知识问答，需要结合知识图谱与RAG。"
        ],
        font_size=15
    )
    add_panel(slide, Inches(6.78), Inches(2.0), Inches(5.98), Inches(4.55), fill=WHITE)
    add_textbox(slide, Inches(7.03), Inches(2.25), Inches(5.4), Inches(0.4), "研究目标", font_size=18, color=PRIMARY, bold=True)
    goals = [
        ("01", "构建离线知识库", "完成多源数据采集、清洗、统一表示与本地索引。"),
        ("02", "实现实体识别与链接", "从文本中识别提及并映射到标准实体。"),
        ("03", "完成跨语言实体对齐", "补充英文侧实体、关系与描述信息。"),
        ("04", "支撑Graph+RAG问答", "为后续知识增强生成和答案溯源打基础。"),
    ]
    y = 2.75
    for idx, (num, title, body) in enumerate(goals):
        box = add_panel(slide, Inches(7.02), Inches(y), Inches(5.45), Inches(0.72), fill=LIGHT_BLUE if idx % 2 == 0 else LIGHT_GREEN)
        add_textbox(slide, Inches(7.18), Inches(y + 0.12), Inches(0.52), Inches(0.3), num, font_size=14, color=ACCENT, bold=True)
        add_textbox(slide, Inches(7.75), Inches(y + 0.10), Inches(1.9), Inches(0.28), title, font_size=14, color=PRIMARY, bold=True)
        add_textbox(slide, Inches(9.4), Inches(y + 0.10), Inches(2.9), Inches(0.34), body, font_size=10.8, color=TEXT)
        y += 0.87

    # Slide 3: Technical route
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "总体技术路线", "围绕“先建库、再训练、后链接、再扩展”的整体流程")
    steps = [
        ("多源数据采集", "Wikidata / 维基百科 / DBP15K", LIGHT_ORANGE, ACCENT),
        ("离线知识库构建", "字段统一、JSON存储、ES索引", LIGHT_BLUE, PRIMARY),
        ("Embedding 与 NER", "Chinese-RoBERTa-wwm-ext-large", RGBColor(244, 239, 255), RGBColor(123, 104, 171)),
        ("实体链接与消歧", "BM25召回 + 向量匹配 + LLM重排", LIGHT_GREEN, GREEN),
        ("跨语言对齐与应用", "LaBSE + GAT + Graph+RAG", WHITE, PRIMARY),
    ]
    left = Inches(0.65)
    top = Inches(2.5)
    box_w = Inches(2.25)
    box_h = Inches(1.55)
    gap = Inches(0.28)
    centers = []
    for idx, (t, b, fill, line) in enumerate(steps):
        l = left + idx * (box_w + gap)
        add_flow_box(slide, l, top, box_w, box_h, t, b, fill, line, title_color=line)
        centers.append((l + box_w / 2, top + box_h / 2))
    for idx in range(len(centers) - 1):
        x1 = centers[idx][0] + box_w / 2 - Inches(0.15)
        x2 = centers[idx + 1][0] - box_w / 2 + Inches(0.15)
        y = centers[idx][1]
        add_arrow(slide, x1, y, x2, y)
    add_panel(slide, Inches(0.8), Inches(4.85), Inches(11.75), Inches(1.25), fill=WHITE)
    add_textbox(slide, Inches(1.0), Inches(5.08), Inches(11.2), Inches(0.55),
                "前一阶段的输出会直接进入后一阶段：离线知识库提供候选与检索基础，Embedding/NER模型提供实体提及识别与统一表示，实体链接完成提及到标准实体的映射，跨语言对齐进一步扩展英文侧知识，最终为Graph+RAG问答提供可检索、可对齐、可溯源的知识支撑。",
                font_size=14, color=TEXT)

    # Slide 4: Offline KG
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "离线知识库构建", "数据来源、预处理流程与当前规模")
    add_panel(slide, Inches(0.6), Inches(2.0), Inches(7.05), Inches(4.8), fill=WHITE)
    fit_image(slide, IMG_PREPROCESS, Inches(0.82), Inches(2.25), Inches(6.6), Inches(3.25))
    add_textbox(slide, Inches(0.95), Inches(5.7), Inches(6.3), Inches(0.28), "图：数据预处理流程", font_size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(7.9), Inches(2.0), Inches(4.85), Inches(4.8), fill=WHITE)
    add_textbox(slide, Inches(8.18), Inches(2.22), Inches(4.3), Inches(0.35), "数据来源与核心统计", font_size=18, color=PRIMARY, bold=True)
    add_bullets(
        slide, Inches(8.2), Inches(2.7), Inches(4.15), Inches(2.1),
        [
            "Wikidata：结构化实体、属性与QID标准标识",
            "中文维基百科：2828页，用于补充描述和上下文",
            "英文维基百科：2844页，用于跨语言扩展",
            "DBP15K：中英文实体对齐标准数据集"
        ],
        font_size=13.5
    )
    add_metric_card(slide, Inches(8.2), Inches(5.05), Inches(1.45), Inches(0.92), "实体规模", "20,000+", accent=PRIMARY)
    add_metric_card(slide, Inches(9.8), Inches(5.05), Inches(1.45), Inches(0.92), "描述填充率", "98.2%", accent=ACCENT)
    add_metric_card(slide, Inches(11.4), Inches(5.05), Inches(1.15), Inches(0.92), "检索响应", "<100ms", accent=GREEN)

    # Slide 5: Embedding & NER
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Embedding 与 NER 模型训练", "多源数据融合、训练配置与阶段结果")
    add_panel(slide, Inches(0.6), Inches(2.0), Inches(4.5), Inches(4.9), fill=WHITE)
    add_textbox(slide, Inches(0.86), Inches(2.2), Inches(3.8), Inches(0.35), "训练数据组成", font_size=18, color=PRIMARY, bold=True)
    add_small_table(
        slide, Inches(0.85), Inches(2.68), Inches(4.0), Inches(2.45),
        ["数据源", "条数"],
        [["自有标注", "3123"], ["CCKS", "1600"], ["MSRA", "50707"], ["自定义补充", "295"], ["合计", "51338"]]
    )
    add_textbox(slide, Inches(0.88), Inches(5.42), Inches(4.0), Inches(0.95),
                "模型采用 Chinese-RoBERTa-wwm-ext-large，\n通过渐进式解冻完成军事领域微调。",
                font_size=12.5, color=TEXT)
    add_panel(slide, Inches(5.35), Inches(2.0), Inches(7.35), Inches(4.9), fill=WHITE)
    add_textbox(slide, Inches(5.65), Inches(2.2), Inches(3.5), Inches(0.35), "训练配置与结果", font_size=18, color=PRIMARY, bold=True)
    add_metric_card(slide, Inches(5.68), Inches(2.75), Inches(1.55), Inches(1.05), "准确率", "98.22%", accent=PRIMARY)
    add_metric_card(slide, Inches(7.36), Inches(2.75), Inches(1.55), Inches(1.05), "精确率", "97.98%", accent=ACCENT)
    add_metric_card(slide, Inches(9.04), Inches(2.75), Inches(1.55), Inches(1.05), "召回率", "98.54%", accent=GREEN)
    add_metric_card(slide, Inches(10.72), Inches(2.75), Inches(1.55), Inches(1.05), "Micro-F1", "98.26%", accent=RGBColor(123, 104, 171))
    add_bullets(
        slide, Inches(5.75), Inches(4.15), Inches(6.3), Inches(1.8),
        [
            "学习率 2×10^-5，batch size 8，最大序列长度 512，训练 5 个 epoch。",
            "NER 作为实体识别入口，为后续实体链接提供提及 span 和统一表示。",
            "当前难点主要集中在简称缩写、嵌套实体和复杂型号表达。"
        ],
        font_size=13.5
    )

    # Slide 6: Entity linking method
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "实体链接方法设计", "候选召回、语义匹配与LLM辅助重排序")
    add_panel(slide, Inches(0.6), Inches(2.0), Inches(6.5), Inches(4.85), fill=WHITE)
    fit_image(slide, IMG_LINKING, Inches(0.85), Inches(2.25), Inches(6.0), Inches(3.15))
    add_textbox(slide, Inches(1.0), Inches(5.65), Inches(5.7), Inches(0.28), "图：实体链接方案", font_size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(7.35), Inches(2.0), Inches(5.35), Inches(4.85), fill=WHITE)
    add_textbox(slide, Inches(7.62), Inches(2.22), Inches(4.8), Inches(0.35), "核心设计", font_size=18, color=PRIMARY, bold=True)
    add_bullets(
        slide, Inches(7.65), Inches(2.72), Inches(4.6), Inches(2.0),
        [
            "候选召回：基于 Elasticsearch 倒排索引和 BM25 多字段检索。",
            "语义匹配：使用与 NER 一致的编码器生成 1024 维向量。",
            "重排序：引入 LLM 处理简称、模糊提及和上下文不足样本。"
        ],
        font_size=14
    )
    add_textbox(slide, Inches(7.62), Inches(4.95), Inches(4.8), Inches(0.32), "实验对比方案", font_size=18, color=PRIMARY, bold=True)
    add_bullets(
        slide, Inches(7.65), Inches(5.38), Inches(4.6), Inches(1.2),
        [
            "纯向量检索 / 纯 ES 文本检索 / 纯 LLM",
            "向量+LLM 重排 / 向量+LLM 动态混合"
        ],
        font_size=13.5
    )

    # Slide 7: Entity linking results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "实体链接实验结果", "以全部数据配置为例")
    add_panel(slide, Inches(0.55), Inches(1.95), Inches(7.15), Inches(4.95), fill=WHITE)
    add_chart_mrr(slide, Inches(0.78), Inches(2.2), Inches(6.6), Inches(3.9))
    add_panel(slide, Inches(7.95), Inches(1.95), Inches(4.8), Inches(4.95), fill=WHITE)
    add_textbox(slide, Inches(8.22), Inches(2.2), Inches(4.2), Inches(0.35), "结果解读", font_size=18, color=PRIMARY, bold=True)
    add_bullets(
        slide, Inches(8.25), Inches(2.72), Inches(4.1), Inches(2.9),
        [
            "纯 ES 文本检索整体最稳定，说明当前任务仍以名称和别名匹配为主。",
            "纯 LLM 在 MRR 和 Hit@1 上略优，适合放在候选重排阶段。",
            "纯向量检索效果较弱，单独承担召回任务时表现不足。",
            "动态混合方案更符合实际使用场景，能够在效果与成本之间取得平衡。"
        ],
        font_size=13.5
    )
    add_metric_card(slide, Inches(8.22), Inches(5.98), Inches(1.35), Inches(0.95), "最佳MRR", "0.7015", accent=ACCENT)
    add_metric_card(slide, Inches(9.72), Inches(5.98), Inches(1.35), Inches(0.95), "最佳Hit@1", "64.19%", accent=PRIMARY)
    add_metric_card(slide, Inches(11.22), Inches(5.98), Inches(1.15), Inches(0.95), "最高Hit@10", "83.78%", accent=GREEN)

    # Slide 8: Alignment method
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "跨语言实体对齐方法", "LaBSE + GAT + 对比学习")
    add_panel(slide, Inches(0.6), Inches(2.0), Inches(6.45), Inches(4.85), fill=WHITE)
    fit_image(slide, IMG_ALIGN, Inches(0.9), Inches(2.25), Inches(5.9), Inches(3.15))
    add_textbox(slide, Inches(1.0), Inches(5.68), Inches(5.7), Inches(0.28), "图：实体对齐训练流程", font_size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(7.3), Inches(2.0), Inches(5.4), Inches(4.85), fill=WHITE)
    add_textbox(slide, Inches(7.58), Inches(2.22), Inches(4.6), Inches(0.35), "方法设计与参数", font_size=18, color=PRIMARY, bold=True)
    add_bullets(
        slide, Inches(7.62), Inches(2.7), Inches(4.6), Inches(2.35),
        [
            "文本编码：LaBSE 生成 768 维跨语言语义表示。",
            "结构建模：单层单头 GAT 聚合一跳邻域信息。",
            "训练目标：NCE Softmax 对比学习，拉近对齐实体、拉远非对齐实体。"
        ],
        font_size=13.5
    )
    add_small_table(
        slide, Inches(7.62), Inches(5.2), Inches(4.4), Inches(1.3),
        ["参数", "取值"],
        [["batch_size", "64"], ["neighbor", "20"], ["lr", "1e-6"], ["dropout", "0.3"]]
    )

    # Slide 9: Alignment results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "跨语言实体对齐结果与作用", "当前模块已具备稳定的中英实体映射能力")
    add_panel(slide, Inches(0.62), Inches(2.0), Inches(4.0), Inches(4.85), fill=WHITE)
    add_textbox(slide, Inches(0.95), Inches(2.25), Inches(3.4), Inches(0.38), "核心指标", font_size=18, color=PRIMARY, bold=True)
    add_metric_card(slide, Inches(0.95), Inches(2.9), Inches(1.0), Inches(1.1), "Hit@1", "0.70", accent=PRIMARY)
    add_metric_card(slide, Inches(2.15), Inches(2.9), Inches(1.0), Inches(1.1), "Hit@10", "0.84", accent=ACCENT)
    add_metric_card(slide, Inches(3.35), Inches(2.9), Inches(1.0), Inches(1.1), "MRR", "0.76", accent=GREEN)
    add_bullets(
        slide, Inches(1.0), Inches(4.45), Inches(3.3), Inches(1.7),
        [
            "文本语义与邻域结构联合建模后，对齐排序更稳定。",
            "轻量 GAT 更适合当前数据规模与训练条件。"
        ],
        font_size=12.5
    )
    add_panel(slide, Inches(4.92), Inches(2.0), Inches(7.75), Inches(4.85), fill=WHITE)
    add_textbox(slide, Inches(5.2), Inches(2.25), Inches(6.9), Inches(0.38), "对项目整体的意义", font_size=18, color=PRIMARY, bold=True)
    impacts = [
        ("知识补全", "把英文侧等价实体、属性与关系补充到现有知识库中。"),
        ("证据扩展", "为后续问答提供更多跨语言描述和外部知识来源。"),
        ("系统衔接", "让实体链接结果可以继续向跨语言检索和Graph+RAG扩展。"),
    ]
    y = 2.95
    for idx, (t, d) in enumerate(impacts):
        fill = LIGHT_BLUE if idx == 0 else LIGHT_ORANGE if idx == 1 else LIGHT_GREEN
        add_panel(slide, Inches(5.22), Inches(y), Inches(6.95), Inches(0.95), fill=fill)
        add_textbox(slide, Inches(5.45), Inches(y + 0.18), Inches(1.5), Inches(0.25), t, font_size=14, color=PRIMARY, bold=True)
        add_textbox(slide, Inches(6.68), Inches(y + 0.15), Inches(5.1), Inches(0.4), d, font_size=12.2, color=TEXT)
        y += 1.12

    # Slide 10: Achievements & innovation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "已完成工作与阶段性成果", "中期阶段已经形成的关键结果")
    cards = [
        ("离线知识库", "已完成多源数据清洗、字段统一和本地索引，形成约 2 万实体知识库。"),
        ("领域 NER", "完成多源数据融合训练，Micro-F1 达到 98.26%。"),
        ("实体链接", "完成 5 种方案对比，明确纯文本检索 + LLM 辅助的有效路径。"),
        ("跨语言对齐", "完成 LaBSE + GAT 模块搭建与评估，具备知识扩展能力。"),
    ]
    positions = [(0.65, 2.15), (6.85, 2.15), (0.65, 4.3), (6.85, 4.3)]
    fills = [LIGHT_ORANGE, LIGHT_BLUE, LIGHT_GREEN, RGBColor(244, 239, 255)]
    for (x, y), (title, body), fill in zip(positions, cards, fills):
        add_panel(slide, Inches(x), Inches(y), Inches(5.75), Inches(1.65), fill=fill)
        add_textbox(slide, Inches(x + 0.22), Inches(y + 0.18), Inches(2.4), Inches(0.3), title, font_size=17, color=PRIMARY, bold=True)
        add_textbox(slide, Inches(x + 0.22), Inches(y + 0.58), Inches(5.15), Inches(0.7), body, font_size=13, color=TEXT)

    # Slide 11: Future plan
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "后续工作计划", "围绕模块优化、系统联调与论文完善推进")
    add_panel(slide, Inches(0.6), Inches(2.0), Inches(5.2), Inches(4.8), fill=WHITE)
    add_textbox(slide, Inches(0.88), Inches(2.22), Inches(4.6), Inches(0.35), "近期重点", font_size=18, color=PRIMARY, bold=True)
    add_bullets(
        slide, Inches(0.92), Inches(2.75), Inches(4.45), Inches(3.6),
        [
            "继续优化实体链接模块，重点处理简称、别名变体和长尾型号样本。",
            "补充跨语言实体对齐实验，观察其在知识补全和问答任务中的实际增益。",
            "推进 Graph+RAG 相关实验与系统联调，验证整体流程的稳定性。"
        ],
        font_size=14
    )
    add_panel(slide, Inches(6.05), Inches(2.0), Inches(6.65), Inches(4.8), fill=WHITE)
    add_textbox(slide, Inches(6.32), Inches(2.22), Inches(4.8), Inches(0.35), "时间安排", font_size=18, color=PRIMARY, bold=True)
    add_timeline(slide, Inches(6.35), Inches(3.0), Inches(5.85), Inches(0.95))
    add_textbox(slide, Inches(6.35), Inches(5.95), Inches(5.7), Inches(0.52),
                "前期主流程已经完成，后续工作以补实验、做联调和整理论文为主，整体时间安排可行。",
                font_size=13, color=TEXT)

    # Slide 12: Thanks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(247, 242, 235))
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(5.8), Inches(13.333), Inches(1.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(1.0), Inches(1.45), Inches(11.2), Inches(0.7), "感谢各位老师的聆听", font_size=28, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(2.3), Inches(11.2), Inches(0.5), "请各位老师批评指正", font_size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(6.2), Inches(11.2), Inches(0.4), "Q & A", font_size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(6.68), Inches(11.2), Inches(0.28), f"{AUTHOR}  |  {MAJOR}  |  {STUDENT_ID}", font_size=12, color=RGBColor(220, 228, 240), align=PP_ALIGN.CENTER)

    prs.save(str(OUT))


if __name__ == "__main__":
    build_presentation()
