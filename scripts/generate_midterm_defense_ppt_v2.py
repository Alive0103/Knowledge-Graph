# -*- coding: utf-8 -*-
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"D:/work/毕设/知识图谱/Knowledge-Graph")
OUT = ROOT / "王越洋-中期答辩PPT.pptx"

IMG_PREPROCESS = Path(r"C:/Users/alive/AppData/Roaming/Typora/typora-user-images/image-20260313131606367.png")
IMG_LINKING = Path(r"C:/Users/alive/AppData/Roaming/Typora/typora-user-images/image-20260313131738574.png")
IMG_ALIGN = Path(r"C:/Users/alive/AppData/Roaming/Typora/typora-user-images/image-20260313131752702.png")

TITLE = "融合内外部数据的知识百科构建与检索增强生成系统"
SUBTITLE = "本科毕业设计中期答辩"
AUTHOR = "王越洋"
STUDENT_ID = "22009200894"
MAJOR = "计算机科学与技术"
SCHOOL = "西安电子科技大学"

BG = RGBColor(247, 245, 241)
PRIMARY = RGBColor(29, 58, 96)
ACCENT = RGBColor(198, 108, 45)
GREEN = RGBColor(91, 136, 95)
TEXT = RGBColor(43, 43, 43)
MUTED = RGBColor(104, 111, 120)
LINE = RGBColor(220, 223, 228)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(235, 242, 249)
LIGHT_ORANGE = RGBColor(249, 239, 230)
LIGHT_GREEN = RGBColor(236, 246, 237)
LIGHT_GOLD = RGBColor(250, 247, 227)
LIGHT_RED = RGBColor(251, 238, 235)


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Microsoft YaHei",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    lines,
    size=14,
    color=TEXT,
    bullet=False,
    line_space=1.1,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        run = p.add_run()
        run.text = f"• {line}" if bullet else line
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, fill=WHITE, line=LINE, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_title(slide, title, subtitle=None, idx=None):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.24))
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.color.rgb = PRIMARY

    add_text(slide, Inches(0.55), Inches(0.55), Inches(8.8), Inches(0.5), title, size=24, color=PRIMARY, bold=True)
    if subtitle:
        add_text(slide, Inches(0.58), Inches(1.02), Inches(10.5), Inches(0.35), subtitle, size=11, color=MUTED)
    if idx is not None:
        add_text(slide, Inches(12.3), Inches(0.58), Inches(0.45), Inches(0.3), f"{idx:02d}", size=12, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)


def add_metric(slide, left, top, width, height, label, value, accent=PRIMARY):
    add_panel(slide, left, top, width, height, fill=WHITE)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.1))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.color.rgb = accent
    add_text(slide, left + Inches(0.14), top + Inches(0.16), width - Inches(0.28), Inches(0.22), label, size=10, color=MUTED, bold=True)
    add_text(slide, left + Inches(0.14), top + Inches(0.42), width - Inches(0.28), Inches(0.36), value, size=22, color=accent, bold=True)


def add_arrow_chevron(slide, left, top, width=Inches(0.42), height=Inches(0.3), color=ACCENT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def fit_image(slide, image_path, left, top, width, height):
    if not image_path.exists():
        add_panel(slide, left, top, width, height, fill=LIGHT_BLUE)
        add_text(slide, left, top + height / 2 - Inches(0.12), width, Inches(0.24), "示意图缺失", size=14, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        return None

    with Image.open(image_path) as img:
        img_w, img_h = img.size

    box_w = width / 914400
    box_h = height / 914400
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h

    if img_ratio > box_ratio:
        final_w = width
        final_h = int(width / img_ratio)
        final_left = left
        final_top = top + int((height - final_h) / 2)
    else:
        final_h = height
        final_w = int(height * img_ratio)
        final_top = top
        final_left = left + int((width - final_w) / 2)

    return slide.shapes.add_picture(str(image_path), final_left, final_top, width=final_w, height=final_h)


def add_result_table(slide, left, top, width, height):
    headers = ["方案", "MRR", "Hit@1", "Hit@10"]
    rows = [
        ("纯向量", "0.0590", "0.0360", "-"),
        ("纯文本检索", "0.6925", "0.6284", "0.8311"),
        ("纯LLM", "0.7015", "0.6419", "-"),
        ("向量+LLM重排", "0.0686", "0.0450", "-"),
        ("向量+LLM动态混合", "0.6892", "0.6239", "0.8333"),
    ]
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    col_widths = [Inches(2.15), Inches(0.95), Inches(1.05), Inches(1.05)]
    for i, col_width in enumerate(col_widths):
        table.columns[i].width = col_width

    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else RGBColor(248, 250, 252)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(10.5)
            run.font.color.rgb = TEXT


def add_bar_group(slide, left, top, width, label, value, max_value, color):
    add_text(slide, left, top, Inches(1.8), Inches(0.25), label, size=11, color=TEXT, bold=True)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(1.85), top + Inches(0.04), width, Inches(0.18))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(229, 233, 238)
    bg.line.color.rgb = RGBColor(229, 233, 238)

    bar_w = int(width * (value / max_value))
    fg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(1.85), top + Inches(0.04), bar_w, Inches(0.18))
    fg.fill.solid()
    fg.fill.fore_color.rgb = color
    fg.line.color.rgb = color

    add_text(slide, left + Inches(5.2), top - Inches(0.01), Inches(0.8), Inches(0.26), f"{value:.4f}", size=10, color=MUTED, bold=True)


def add_timeline(slide, left, top):
    items = [
        ("3月下旬-4月中旬", "继续优化实体链接，补误差分析与消融对比"),
        ("4月中下旬", "完成跨语言对齐补实验，并推进 Graph+RAG 联调"),
        ("5月", "整理图表与论文，完成系统展示、修改与定稿"),
    ]
    for idx, (period, task) in enumerate(items):
        y = top + Inches(idx * 1.05)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, y + Inches(0.18), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT if idx == 0 else PRIMARY if idx == 1 else GREEN
        dot.line.color.rgb = dot.fill.fore_color.rgb
        if idx < len(items) - 1:
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(0.065), y + Inches(0.35), Inches(0.03), Inches(0.78))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(210, 215, 222)
            line.line.color.rgb = RGBColor(210, 215, 222)
        add_text(slide, left + Inches(0.32), y, Inches(1.8), Inches(0.24), period, size=11, color=ACCENT if idx == 0 else PRIMARY if idx == 1 else GREEN, bold=True)
        add_panel(slide, left + Inches(2.0), y - Inches(0.03), Inches(4.35), Inches(0.56), fill=WHITE)
        add_text(slide, left + Inches(2.2), y + Inches(0.08), Inches(3.95), Inches(0.3), task, size=11.5, color=TEXT)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(246, 242, 236))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(5.85), Inches(13.333), Inches(1.65))
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.color.rgb = PRIMARY
    add_text(slide, Inches(0.95), Inches(1.1), Inches(11.5), Inches(1.0), TITLE, size=28, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4.35), Inches(2.2), Inches(4.6), Inches(0.38), SUBTITLE, size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(4.18), Inches(3.05), Inches(5.0), Inches(1.3), fill=WHITE)
    add_text(slide, Inches(4.45), Inches(3.28), Inches(1.15), Inches(0.28), "姓名", size=13, color=MUTED, bold=True)
    add_text(slide, Inches(5.3), Inches(3.28), Inches(1.2), Inches(0.28), AUTHOR, size=14, color=TEXT, bold=True)
    add_text(slide, Inches(6.45), Inches(3.28), Inches(1.15), Inches(0.28), "学号", size=13, color=MUTED, bold=True)
    add_text(slide, Inches(7.28), Inches(3.28), Inches(1.4), Inches(0.28), STUDENT_ID, size=14, color=TEXT, bold=True)
    add_text(slide, Inches(4.45), Inches(3.72), Inches(1.15), Inches(0.28), "专业", size=13, color=MUTED, bold=True)
    add_text(slide, Inches(5.3), Inches(3.72), Inches(2.3), Inches(0.28), MAJOR, size=14, color=TEXT, bold=True)
    add_text(slide, Inches(6.45), Inches(3.72), Inches(1.15), Inches(0.28), "学校", size=13, color=MUTED, bold=True)
    add_text(slide, Inches(7.28), Inches(3.72), Inches(2.2), Inches(0.28), SCHOOL, size=14, color=TEXT, bold=True)
    add_text(slide, Inches(1.0), Inches(6.32), Inches(11.2), Inches(0.35), "Knowledge Graph · Entity Linking · Cross-lingual Alignment · Graph+RAG", size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 2. Background
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "研究背景与课题目标", "围绕知识构建、实体理解与问答增强三个层面展开", 2)
    add_panel(slide, Inches(0.62), Inches(1.7), Inches(6.0), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(0.92), Inches(2.0), Inches(2.0), Inches(0.3), "为什么要做", size=18, color=PRIMARY, bold=True)
    add_paragraphs(
        slide,
        Inches(0.95),
        Inches(2.45),
        Inches(5.2),
        Inches(3.4),
        [
            "组织内部知识与外部开放知识长期分散，难以形成统一、可检索、可追溯的知识底座。",
            "仅依赖文档级 RAG 容易出现实体混淆、证据碎片化和答案来源不清晰的问题。",
            "军事领域实体名称复杂，存在简称、别名、型号变体和跨语言差异，更需要实体级知识组织。 ",
        ],
        size=14,
        color=TEXT,
        bullet=True,
    )
    add_panel(slide, Inches(6.95), Inches(1.7), Inches(5.75), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(7.25), Inches(2.0), Inches(2.6), Inches(0.3), "本课题目标", size=18, color=PRIMARY, bold=True)
    goals = [
        ("知识底座", "融合 Wikidata、维基百科与 DBP15K，构建离线可检索知识库。", LIGHT_BLUE),
        ("实体理解", "完成 NER、实体向量表示与实体链接，实现文本到标准实体的映射。", LIGHT_ORANGE),
        ("问答增强", "通过跨语言对齐补充英文侧知识，为 Graph+RAG 提供更完整证据。", LIGHT_GREEN),
    ]
    y = 2.45
    for title, body, fill in goals:
        add_panel(slide, Inches(7.28), Inches(y), Inches(5.1), Inches(0.95), fill=fill)
        add_text(slide, Inches(7.5), Inches(y + 0.15), Inches(1.1), Inches(0.22), title, size=14, color=PRIMARY, bold=True)
        add_text(slide, Inches(8.45), Inches(y + 0.13), Inches(3.65), Inches(0.42), body, size=12.2, color=TEXT)
        y += 1.18

    # 3. Route
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "整体技术路线", "由知识准备到实体理解，再到跨语言扩展与问答增强，形成递进式链路", 3)
    add_text(
        slide,
        Inches(0.72),
        Inches(1.68),
        Inches(11.8),
        Inches(0.45),
        "本课题先完成离线知识库构建，保证实体、属性和文本描述能够统一存储与检索；在此基础上训练领域 NER 与实体表示模型，为后续提及识别和候选匹配提供输入；随后结合文本检索、向量相似度和 LLM 重排序完成实体链接；再通过跨语言实体对齐把英文侧等价实体及其关系补入知识库，最终为 Graph+RAG 问答提供更完整、更可追溯的知识支撑。",
        size=12.3,
        color=TEXT,
    )
    stages = [
        ("阶段一", "离线知识库构建", "Wikidata + 维基百科 + DBP15K\n清洗、补全、统一表示与 ES 索引", LIGHT_BLUE),
        ("阶段二", "Embedding / NER 训练", "Chinese-RoBERTa-wwm-ext-large\n提及识别与实体语义表示", LIGHT_ORANGE),
        ("阶段三", "实体链接", "BM25 多字段召回 + 向量相似度\nLLM 辅助重排序与消歧", LIGHT_GREEN),
        ("阶段四", "跨语言实体对齐", "LaBSE + GAT + 对比学习\n中英等价实体映射", LIGHT_GOLD),
        ("阶段五", "Graph+RAG 验证", "知识检索、证据组织与答案生成\n比较准确性、完整性和可追溯性", LIGHT_RED),
    ]
    x = 0.68
    for idx, (tag, title, body, fill) in enumerate(stages):
        add_panel(slide, Inches(x), Inches(2.45), Inches(2.22), Inches(2.8), fill=fill)
        add_text(slide, Inches(x + 0.18), Inches(2.64), Inches(0.8), Inches(0.2), tag, size=11, color=ACCENT, bold=True)
        add_text(slide, Inches(x + 0.18), Inches(2.96), Inches(1.85), Inches(0.48), title, size=16, color=PRIMARY, bold=True)
        add_text(slide, Inches(x + 0.18), Inches(3.68), Inches(1.85), Inches(1.15), body, size=11.5, color=TEXT)
        if idx < len(stages) - 1:
            add_arrow_chevron(slide, Inches(x + 2.3), Inches(3.6))
        x += 2.52
    add_text(slide, Inches(0.74), Inches(5.65), Inches(11.8), Inches(0.35), "递进关系：前一阶段的结构化结果和表示结果，直接作为后一阶段的输入或检索基础，整体上是一条可落地、可扩展的技术链路。", size=12, color=MUTED)

    # 4. KB
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "离线知识库构建", "完成多源数据汇聚、字段补全与本地检索能力建设", 4)
    add_panel(slide, Inches(0.62), Inches(1.8), Inches(5.8), Inches(4.9), fill=WHITE)
    fit_image(slide, IMG_PREPROCESS, Inches(0.88), Inches(2.1), Inches(5.25), Inches(2.9))
    add_text(slide, Inches(1.15), Inches(5.22), Inches(4.7), Inches(0.24), "数据处理链路示意：源数据清洗、内容补全、实体抽取增强与索引存储。", size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(6.72), Inches(1.8), Inches(6.0), Inches(4.9), fill=WHITE)
    add_text(slide, Inches(7.0), Inches(2.05), Inches(1.8), Inches(0.28), "数据来源与处理", size=18, color=PRIMARY, bold=True)
    add_paragraphs(
        slide,
        Inches(7.02),
        Inches(2.45),
        Inches(5.35),
        Inches(1.55),
        [
            "结构化知识以 Wikidata 为主，补充中文维基 2828 页与英文维基 2844 页文本。",
            "引入 DBP15K 的中英实体对和三元组，为后续跨语言实体对齐提供标准数据基础。",
            "统一整理为 JSON，并建立 Elasticsearch 多字段检索与向量索引。 ",
        ],
        size=12.6,
        bullet=True,
    )
    add_metric(slide, Inches(7.02), Inches(4.5), Inches(1.25), Inches(0.95), "实体规模", "约2万", PRIMARY)
    add_metric(slide, Inches(8.45), Inches(4.5), Inches(1.35), Inches(0.95), "描述填充率", "98.2%", ACCENT)
    add_metric(slide, Inches(10.0), Inches(4.5), Inches(1.35), Inches(0.95), "别名填充率", "95.7%", GREEN)
    add_metric(slide, Inches(11.55), Inches(4.5), Inches(1.0), Inches(0.95), "响应", "<100ms", PRIMARY)
    add_text(slide, Inches(7.02), Inches(5.78), Inches(5.4), Inches(0.42), "阶段结果：知识库已经能够支撑实体名/别名全文检索、向量语义检索以及后续候选召回实验。", size=12.3, color=TEXT)

    # 5. NER
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Embedding 与 NER 模型训练", "以 Chinese-RoBERTa-wwm-ext-large 作为统一编码器", 5)
    add_panel(slide, Inches(0.62), Inches(1.8), Inches(6.0), Inches(4.9), fill=WHITE)
    add_text(slide, Inches(0.92), Inches(2.08), Inches(2.6), Inches(0.28), "训练思路", size=18, color=PRIMARY, bold=True)
    blocks = [
        ("输入数据", "领域实体标注样本 + 知识库实体文本", LIGHT_BLUE),
        ("统一编码器", "Chinese-RoBERTa-wwm-ext-large\n1024 维隐藏表示", LIGHT_ORANGE),
        ("下游输出", "NER 提及识别 + 实体向量表示", LIGHT_GREEN),
    ]
    y = 2.6
    for idx, (title, body, fill) in enumerate(blocks):
        add_panel(slide, Inches(1.08), Inches(y), Inches(4.95), Inches(0.86), fill=fill)
        add_text(slide, Inches(1.3), Inches(y + 0.12), Inches(1.35), Inches(0.2), title, size=13, color=PRIMARY, bold=True)
        add_text(slide, Inches(2.45), Inches(y + 0.1), Inches(3.2), Inches(0.36), body, size=12, color=TEXT)
        if idx < len(blocks) - 1:
            add_arrow_chevron(slide, Inches(3.12), Inches(y + 0.93), Inches(0.42), Inches(0.26), PRIMARY)
        y += 1.2
    add_text(slide, Inches(0.92), Inches(6.0), Inches(5.4), Inches(0.35), "这样做的原因是：NER 负责稳定识别提及边界，向量表示则为知识库实体匹配和后续实体链接提供统一语义空间。", size=12.2, color=MUTED)
    add_panel(slide, Inches(6.92), Inches(1.8), Inches(5.8), Inches(4.9), fill=WHITE)
    add_text(slide, Inches(7.22), Inches(2.08), Inches(2.8), Inches(0.28), "当前识别效果", size=18, color=PRIMARY, bold=True)
    add_metric(slide, Inches(7.22), Inches(2.7), Inches(1.2), Inches(1.0), "准确率", "98.22%", PRIMARY)
    add_metric(slide, Inches(8.62), Inches(2.7), Inches(1.2), Inches(1.0), "精确率", "97.98%", ACCENT)
    add_metric(slide, Inches(10.02), Inches(2.7), Inches(1.2), Inches(1.0), "召回率", "98.54%", GREEN)
    add_metric(slide, Inches(11.42), Inches(2.7), Inches(1.0), Inches(1.0), "F1", "98.26%", PRIMARY)
    add_paragraphs(
        slide,
        Inches(7.22),
        Inches(4.15),
        Inches(4.95),
        Inches(1.4),
        [
            "当前 NER 已能较稳定识别人名、地名、组织以及军事装备、型号、事件等专有实体。",
            "这一模块的输出会直接作为实体链接阶段的提及输入，并反向支撑误差分析。 ",
        ],
        size=12.8,
        bullet=True,
    )

    # 6. Entity linking method
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "实体链接方法设计", "将文本中的实体提及映射到知识库标准实体", 6)
    add_panel(slide, Inches(0.62), Inches(1.8), Inches(6.0), Inches(4.9), fill=WHITE)
    fit_image(slide, IMG_LINKING, Inches(0.85), Inches(2.02), Inches(5.45), Inches(3.3))
    add_text(slide, Inches(1.05), Inches(5.48), Inches(5.0), Inches(0.24), "实体链接流程：候选召回、候选筛选、LLM 辅助消歧和最终输出。", size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(6.92), Inches(1.8), Inches(5.8), Inches(4.9), fill=WHITE)
    add_text(slide, Inches(7.22), Inches(2.08), Inches(2.2), Inches(0.28), "方法组成", size=18, color=PRIMARY, bold=True)
    add_paragraphs(
        slide,
        Inches(7.25),
        Inches(2.45),
        Inches(5.0),
        Inches(2.2),
        [
            "文本侧：Elasticsearch 多字段检索，利用 label、alias、description 等字段做 BM25 候选召回。",
            "语义侧：使用与 NER 一致的 Chinese-RoBERTa 编码提及与实体，做向量相似度排序。",
            "判别侧：用 LLM 对候选进行辅助重排序，重点处理简称、模糊提及和上下文不足样本。 ",
        ],
        size=12.7,
        bullet=True,
    )
    add_panel(slide, Inches(7.25), Inches(5.0), Inches(4.95), Inches(0.95), fill=LIGHT_ORANGE)
    add_text(slide, Inches(7.48), Inches(5.18), Inches(1.1), Inches(0.22), "设计重点", size=13, color=PRIMARY, bold=True)
    add_text(slide, Inches(8.45), Inches(5.12), Inches(3.45), Inches(0.34), "不是让 LLM 直接替代检索，而是放在候选集合上做高价值消歧，平衡效果与成本。", size=11.8, color=TEXT)

    # 7. Entity linking result
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "实体链接实验结果", "重点观察不同方案在召回与排序上的表现差异", 7)
    add_panel(slide, Inches(0.62), Inches(1.8), Inches(6.2), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(0.92), Inches(2.06), Inches(2.4), Inches(0.28), "主要结果表", size=18, color=PRIMARY, bold=True)
    add_result_table(slide, Inches(0.92), Inches(2.55), Inches(5.55), Inches(2.65))
    add_text(slide, Inches(0.95), Inches(5.5), Inches(5.45), Inches(0.35), "说明：纯文本检索已较稳定；纯 LLM 在 MRR 和 Hit@1 上略优，但成本更高；动态混合方案更适合工程使用。", size=11.5, color=MUTED)
    add_panel(slide, Inches(7.05), Inches(1.8), Inches(5.65), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(7.35), Inches(2.06), Inches(2.6), Inches(0.28), "MRR 对比", size=18, color=PRIMARY, bold=True)
    add_bar_group(slide, Inches(7.35), Inches(2.75), Inches(3.2), "纯向量", 0.0590, 0.75, RGBColor(198, 205, 216))
    add_bar_group(slide, Inches(7.35), Inches(3.22), Inches(3.2), "纯文本检索", 0.6925, 0.75, PRIMARY)
    add_bar_group(slide, Inches(7.35), Inches(3.69), Inches(3.2), "纯LLM", 0.7015, 0.75, ACCENT)
    add_bar_group(slide, Inches(7.35), Inches(4.16), Inches(3.2), "向量+LLM重排", 0.0686, 0.75, RGBColor(170, 177, 188))
    add_bar_group(slide, Inches(7.35), Inches(4.63), Inches(3.2), "动态混合", 0.6892, 0.75, GREEN)
    add_metric(slide, Inches(7.38), Inches(5.45), Inches(1.45), Inches(0.95), "最佳 MRR", "0.7015", ACCENT)
    add_metric(slide, Inches(9.06), Inches(5.45), Inches(1.45), Inches(0.95), "最佳 Hit@1", "64.19%", PRIMARY)
    add_metric(slide, Inches(10.74), Inches(5.45), Inches(1.55), Inches(0.95), "最高 Hit@10", "83.78%", GREEN)

    # 8. Alignment method
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "跨语言实体对齐方法", "使用文本语义与邻域结构联合建模中英等价实体", 8)
    add_panel(slide, Inches(0.62), Inches(1.8), Inches(5.9), Inches(4.95), fill=WHITE)
    fit_image(slide, IMG_ALIGN, Inches(0.95), Inches(2.02), Inches(5.2), Inches(3.65))
    add_text(slide, Inches(1.05), Inches(5.9), Inches(5.0), Inches(0.24), "对齐流程：LaBSE 编码文本，GAT 聚合一跳邻域，再用对比学习优化中英实体映射。", size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(6.82), Inches(1.8), Inches(5.9), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(7.12), Inches(2.06), Inches(2.9), Inches(0.28), "关键方案", size=18, color=PRIMARY, bold=True)
    add_paragraphs(
        slide,
        Inches(7.15),
        Inches(2.45),
        Inches(5.0),
        Inches(2.25),
        [
            "数据基础：选取 DBP15K 的 zh_en 子集，并结合维基跨语言链接信息做预处理。",
            "文本编码：使用 LaBSE 生成 768 维跨语言语义表示。",
            "结构聚合：采用单层单头 GAT 编码实体的一跳邻域结构。",
            "训练目标：通过对比学习拉近等价实体、拉远非对齐实体。 ",
        ],
        size=12.4,
        bullet=True,
    )
    add_panel(slide, Inches(7.2), Inches(5.08), Inches(4.95), Inches(0.95), fill=LIGHT_GOLD)
    add_text(slide, Inches(7.45), Inches(5.26), Inches(1.05), Inches(0.22), "输出作用", size=13, color=PRIMARY, bold=True)
    add_text(slide, Inches(8.36), Inches(5.2), Inches(3.45), Inches(0.34), "将中文知识库中的实体连接到英文侧等价实体，支持知识补全与证据扩展。", size=11.8, color=TEXT)

    # 9. Alignment result
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "跨语言实体对齐结果与意义", "当前模块已具备较稳定的中英实体映射能力", 9)
    add_panel(slide, Inches(0.62), Inches(1.8), Inches(4.1), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(0.95), Inches(2.08), Inches(2.3), Inches(0.28), "核心指标", size=18, color=PRIMARY, bold=True)
    add_metric(slide, Inches(0.98), Inches(2.75), Inches(1.0), Inches(1.02), "Hit@1", "0.70", PRIMARY)
    add_metric(slide, Inches(2.16), Inches(2.75), Inches(1.0), Inches(1.02), "Hit@10", "0.84", ACCENT)
    add_metric(slide, Inches(3.34), Inches(2.75), Inches(1.0), Inches(1.02), "MRR", "0.76", GREEN)
    add_paragraphs(
        slide,
        Inches(0.98),
        Inches(4.28),
        Inches(3.3),
        Inches(1.45),
        [
            "文本语义与邻域结构联合建模后，对齐排序更稳，能较好地区分等价实体。",
            "模块链路已跑通，后续重点放在补实验和知识补全场景验证。 ",
        ],
        size=12.2,
        bullet=True,
    )
    add_panel(slide, Inches(5.0), Inches(1.8), Inches(7.7), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(5.32), Inches(2.08), Inches(2.9), Inches(0.28), "对整体系统的作用", size=18, color=PRIMARY, bold=True)
    cards = [
        ("知识补全", "把英文侧等价实体、属性和关系补充到现有中文知识库中。", LIGHT_BLUE),
        ("证据扩展", "在后续问答阶段引入更多跨语言描述、外部来源和可验证证据。", LIGHT_ORANGE),
        ("系统衔接", "让实体链接结果继续扩展到跨语言检索和 Graph+RAG 验证阶段。", LIGHT_GREEN),
    ]
    y = 2.7
    for title, body, fill in cards:
        add_panel(slide, Inches(5.35), Inches(y), Inches(6.9), Inches(0.92), fill=fill)
        add_text(slide, Inches(5.62), Inches(y + 0.16), Inches(1.05), Inches(0.22), title, size=14, color=PRIMARY, bold=True)
        add_text(slide, Inches(6.7), Inches(y + 0.13), Inches(5.2), Inches(0.34), body, size=12.3, color=TEXT)
        y += 1.13

    # 10. Achievements
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "已完成工作与阶段性成果", "中期阶段已经形成从知识准备到实体理解的主流程", 10)
    achievements = [
        ("离线知识库", "完成多源数据清洗、字段统一、描述与别名补全，并建立本地 ES 索引。", LIGHT_BLUE),
        ("领域 NER", "完成编码器训练与评估，实体识别效果较稳定，可支撑后续链接输入。", LIGHT_ORANGE),
        ("实体链接", "完成 5 种方案对比，明确文本检索与 LLM 辅助消歧的有效组合。", LIGHT_GREEN),
        ("跨语言对齐", "完成 LaBSE + GAT 方法链路与基本评估，具备知识扩展能力。", LIGHT_GOLD),
    ]
    positions = [(0.72, 2.0), (6.92, 2.0), (0.72, 4.15), (6.92, 4.15)]
    for (x, y), (title, body, fill) in zip(positions, achievements):
        add_panel(slide, Inches(x), Inches(y), Inches(5.65), Inches(1.6), fill=fill)
        add_text(slide, Inches(x + 0.24), Inches(y + 0.18), Inches(1.8), Inches(0.26), title, size=17, color=PRIMARY, bold=True)
        add_text(slide, Inches(x + 0.24), Inches(y + 0.56), Inches(4.95), Inches(0.62), body, size=13, color=TEXT)
    add_text(slide, Inches(0.85), Inches(6.22), Inches(11.7), Inches(0.34), "当前判断：核心研究链路已经基本贯通，后续工作的重点不再是从零搭框架，而是围绕模块优化、补实验、系统联调和论文完善继续推进。", size=12.4, color=MUTED)

    # 11. Future plan
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "后续工作与进度安排", "围绕模块优化、系统联调和论文收束继续推进", 11)
    add_panel(slide, Inches(0.62), Inches(1.8), Inches(5.7), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(0.95), Inches(2.08), Inches(2.4), Inches(0.28), "后续重点", size=18, color=PRIMARY, bold=True)
    add_text(
        slide,
        Inches(0.98),
        Inches(2.48),
        Inches(5.0),
        Inches(3.7),
        "后续工作会继续围绕实体链接、跨语言对齐和 Graph+RAG 系统验证三部分展开。实体链接方面，重点处理简称、别名变体、长尾型号等复杂样本，进一步补充候选实体信息并完善误差分析；跨语言对齐方面，继续补充标准化对比实验和参数分析，验证对齐结果在知识补全中的实际价值；系统集成方面，将实体链接结果和跨语言对齐结果进一步接入问答流程，比较不同检索与生成组织方式对答案准确性、完整性和可追溯性的影响。",
        size=12.6,
        color=TEXT,
    )
    add_text(slide, Inches(0.98), Inches(6.05), Inches(5.0), Inches(0.3), "整体上，当前前置模块已经具备较好的承接关系，因此后续安排以优化、补实验和集成为主，具备较强可行性。", size=12.2, color=MUTED)
    add_panel(slide, Inches(6.55), Inches(1.8), Inches(6.15), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(6.88), Inches(2.08), Inches(2.4), Inches(0.28), "时间安排", size=18, color=PRIMARY, bold=True)
    add_timeline(slide, Inches(6.95), Inches(2.72))

    # 12. Thanks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(247, 242, 235))
    bottom = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(5.85), Inches(13.333), Inches(1.65))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = PRIMARY
    bottom.line.color.rgb = PRIMARY
    add_text(slide, Inches(1.0), Inches(1.35), Inches(11.2), Inches(0.55), "感谢各位老师的聆听", size=28, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.0), Inches(2.2), Inches(11.2), Inches(0.35), "恳请批评指正", size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.0), Inches(6.28), Inches(11.2), Inches(0.38), "Q & A", size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.0), Inches(6.7), Inches(11.2), Inches(0.24), f"{AUTHOR}  |  {MAJOR}  |  {STUDENT_ID}", size=12, color=RGBColor(226, 232, 241), align=PP_ALIGN.CENTER)

    prs.save(str(OUT))


if __name__ == "__main__":
    build_presentation()
